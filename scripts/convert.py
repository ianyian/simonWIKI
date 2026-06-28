#!/usr/bin/env python3
"""
convert.py — extract text from a document into Markdown for ingestion.

Usage:
    python scripts/convert.py raw/path/to/file.docx
    python scripts/convert.py raw/path/to/file.pdf  --out raw/path/to/file.md

Supports: .docx .pptx .xlsx .pdf .html .htm  (and passes through .md/.txt).
Needs the add-ons from tools/prerequisites.md:
    pip install python-docx python-pptx pdfplumber pypdf openpyxl
Prints the Markdown to stdout, and also writes a .md next to the source unless --no-write.
"""
import sys
import os
import argparse


def convert_docx(path):
    import docx  # python-docx
    d = docx.Document(path)
    return "\n\n".join(p.text for p in d.paragraphs if p.text.strip())


def convert_pptx(path):
    from pptx import Presentation
    out = []
    for i, slide in enumerate(Presentation(path).slides, 1):
        out.append(f"## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                out.append(shape.text_frame.text)
    return "\n\n".join(out)


def convert_xlsx(path):
    from openpyxl import load_workbook
    wb = load_workbook(path, data_only=True)
    out = []
    for ws in wb.worksheets:
        out.append(f"## Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [("" if c is None else str(c)) for c in row]
            if any(cells):
                out.append("| " + " | ".join(cells) + " |")
    return "\n".join(out)


def convert_pdf(path):
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            return "\n\n".join((pg.extract_text() or "") for pg in pdf.pages)
    except Exception:
        from pypdf import PdfReader
        return "\n\n".join((pg.extract_text() or "") for pg in PdfReader(path).pages)


def convert_html(path):
    import re
    with open(path, encoding="utf-8", errors="ignore") as f:
        html = f.read()
    text = re.sub(r"(?is)<(script|style).*?>.*?</\1>", "", html)
    text = re.sub(r"(?s)<[^>]+>", " ", text)
    return re.sub(r"[ \t]+\n", "\n", text).strip()


CONVERTERS = {
    ".docx": convert_docx, ".pptx": convert_pptx, ".xlsx": convert_xlsx,
    ".pdf": convert_pdf, ".html": convert_html, ".htm": convert_html,
}


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("path")
    ap.add_argument("--out")
    ap.add_argument("--no-write", action="store_true")
    args = ap.parse_args()

    ext = os.path.splitext(args.path)[1].lower()
    if ext in (".md", ".txt"):
        with open(args.path, encoding="utf-8", errors="ignore") as f:
            md = f.read()
    elif ext in CONVERTERS:
        try:
            md = CONVERTERS[ext](args.path)
        except ImportError as e:
            sys.exit(f"Missing add-on for {ext}: {e}. See tools/prerequisites.md.")
    else:
        sys.exit(f"Unsupported type: {ext}")

    if not args.no_write:
        out = args.out or os.path.splitext(args.path)[0] + ".md"
        with open(out, "w", encoding="utf-8") as f:
            f.write(md)
        print(f"# wrote {out}", file=sys.stderr)
    print(md)


if __name__ == "__main__":
    main()
